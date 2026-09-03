"""Multimodal document extraction for RAG.

Reliably supported when parsers succeed:
- TXT, CSV, PDF (text + tables + OCR fallback for scanned pages)
- DOCX (paragraphs + tables)
- XLSX (sheets as tables)
- PPTX (slide text + table shapes)
- PNG/JPG/WEBP (OCR / vision when available)

Graph/chart values are labeled as exact, estimated, or interpretation.
Unsupported or empty extracts are reported honestly — never invented.
"""

from __future__ import annotations

import io
import os
import re
from typing import Any

MAX_FILE_BYTES = 25 * 1024 * 1024

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".txt",
    ".csv",
    ".xlsx",
    ".pptx",
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
}

MIME_BY_EXT = {
    ".pdf": "application/pdf",
    ".docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    ".txt": "text/plain",
    ".csv": "text/csv",
    ".xlsx": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    ".pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    ".png": "image/png",
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".webp": "image/webp",
}


class DocumentValidationError(ValueError):
    pass


def validate_upload(filename: str, data: bytes) -> str:
    ext = os.path.splitext(filename or "")[1].lower()
    if ext not in SUPPORTED_EXTENSIONS:
        raise DocumentValidationError(
            f"Unsupported file type '{ext}'. Supported: {', '.join(sorted(SUPPORTED_EXTENSIONS))}"
        )
    if not data:
        raise DocumentValidationError("Empty file.")
    if len(data) > MAX_FILE_BYTES:
        raise DocumentValidationError(f"File exceeds {MAX_FILE_BYTES // (1024 * 1024)} MB limit.")
    return ext


def process_document(filename: str, data: bytes) -> dict[str, Any]:
    ext = validate_upload(filename, data)
    notes: list[str] = []
    tables: list[dict[str, Any]] = []
    visuals: list[dict[str, Any]] = []
    text_parts: list[str] = []
    ocr_used = False

    try:
        if ext == ".txt":
            text_parts.append(data.decode("utf-8", errors="ignore"))
        elif ext == ".csv":
            t, tbls = _extract_csv(data)
            text_parts.append(t)
            tables.extend(tbls)
        elif ext == ".xlsx":
            t, tbls = _extract_xlsx(data)
            text_parts.append(t)
            tables.extend(tbls)
        elif ext == ".docx":
            t, tbls = _extract_docx(data)
            text_parts.append(t)
            tables.extend(tbls)
        elif ext == ".pptx":
            t, tbls = _extract_pptx(data)
            text_parts.append(t)
            tables.extend(tbls)
        elif ext == ".pdf":
            t, tbls, vis, ocr_used, extra_notes = _extract_pdf(data)
            text_parts.append(t)
            tables.extend(tbls)
            visuals.extend(vis)
            notes.extend(extra_notes)
        elif ext in {".png", ".jpg", ".jpeg", ".webp"}:
            t, vis, ocr_used, extra_notes = _extract_image(data, filename)
            text_parts.append(t)
            visuals.extend(vis)
            notes.extend(extra_notes)
    except DocumentValidationError:
        raise
    except Exception as exc:
        notes.append(f"Parser error for {filename}: {exc}")

    full_text = "\n\n".join(p.strip() for p in text_parts if p and p.strip())
    table_md = "\n\n".join(tbl.get("markdown", "") for tbl in tables if tbl.get("markdown"))
    visual_md = "\n\n".join(_visual_block(v) for v in visuals)

    indexed = "\n\n".join(x for x in (full_text, table_md, visual_md) if x).strip()
    if not indexed:
        notes.append("No extractable text, tables, or visual content was found.")
        indexed = ""

    return {
        "filename": filename,
        "extension": ext,
        "mime": MIME_BY_EXT.get(ext, "application/octet-stream"),
        "size_bytes": len(data),
        "text": full_text,
        "tables": tables,
        "visuals": visuals,
        "ocr_used": ocr_used,
        "parse_notes": notes,
        "indexed_text": indexed,
    }


def _visual_block(v: dict[str, Any]) -> str:
    kind = v.get("value_kind", "interpretation")
    loc = v.get("location", "visual")
    desc = v.get("description", "")
    return (
        f"[Visual:{loc}] value_kind={kind}\n{desc}"
    )


def _table_to_markdown(headers: list[str], rows: list[list[str]]) -> str:
    if not headers and not rows:
        return ""
    headers = [str(h or "").strip() or f"Col{i+1}" for i, h in enumerate(headers)]
    if not headers and rows:
        headers = [f"Col{i+1}" for i in range(len(rows[0]))]
    line = "| " + " | ".join(headers) + " |"
    sep = "| " + " | ".join("---" for _ in headers) + " |"
    body = []
    for row in rows:
        cells = [str(c if c is not None else "").strip() for c in row]
        while len(cells) < len(headers):
            cells.append("")
        body.append("| " + " | ".join(cells[: len(headers)]) + " |")
    return "\n".join([line, sep, *body])


def _extract_csv(data: bytes) -> tuple[str, list[dict[str, Any]]]:
    import pandas as pd

    df = pd.read_csv(io.BytesIO(data))
    headers = [str(c) for c in df.columns]
    rows = df.astype(str).values.tolist()
    md = _table_to_markdown(headers, rows[:200])
    return md, [{"headers": headers, "rows": rows[:200], "markdown": md, "source": "csv"}]


def _extract_xlsx(data: bytes) -> tuple[str, list[dict[str, Any]]]:
    import pandas as pd

    tables = []
    parts = []
    xl = pd.ExcelFile(io.BytesIO(data))
    for sheet in xl.sheet_names:
        df = xl.parse(sheet)
        headers = [str(c) for c in df.columns]
        rows = df.astype(str).fillna("").values.tolist()
        md = f"### Sheet: {sheet}\n" + _table_to_markdown(headers, rows[:200])
        tables.append({"headers": headers, "rows": rows[:200], "markdown": md, "source": f"xlsx:{sheet}"})
        parts.append(md)
    return "\n\n".join(parts), tables


def _extract_docx(data: bytes) -> tuple[str, list[dict[str, Any]]]:
    tables: list[dict[str, Any]] = []
    paras: list[str] = []
    try:
        from docx import Document

        doc = Document(io.BytesIO(data))
        paras = [p.text.strip() for p in doc.paragraphs if p.text and p.text.strip()]
        for i, tbl in enumerate(doc.tables, start=1):
            grid = [[(cell.text or "").strip() for cell in row.cells] for row in tbl.rows]
            headers = grid[0] if grid else []
            rows = grid[1:] if len(grid) > 1 else []
            md = _table_to_markdown(headers, rows)
            tables.append({"headers": headers, "rows": rows, "markdown": md, "source": f"docx:table-{i}"})
    except Exception:
        try:
            import docx2txt

            paras = [docx2txt.process(io.BytesIO(data))]
        except Exception:
            paras = []
    body = "\n".join(paras)
    extra = "\n\n".join(t["markdown"] for t in tables)
    return (body + ("\n\n" + extra if extra else "")).strip(), tables


def _extract_pptx(data: bytes) -> tuple[str, list[dict[str, Any]]]:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE

    prs = Presentation(io.BytesIO(data))
    parts = []
    tables = []
    for idx, slide in enumerate(prs.slides, start=1):
        bits = [f"## Slide {idx}"]
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                bits.append(shape.text.strip())
            if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                tbl = shape.table
                grid = [[cell.text.strip() for cell in row.cells] for row in tbl.rows]
                headers = grid[0] if grid else []
                rows = grid[1:] if len(grid) > 1 else []
                md = _table_to_markdown(headers, rows)
                tables.append({"headers": headers, "rows": rows, "markdown": md, "source": f"pptx:slide-{idx}"})
                bits.append(md)
        parts.append("\n".join(b for b in bits if b))
    return "\n\n".join(parts), tables


def _extract_pdf(data: bytes) -> tuple[str, list[dict[str, Any]], list[dict[str, Any]], bool, list[str]]:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    page_texts = []
    notes = []
    ocr_used = False
    visuals: list[dict[str, Any]] = []
    tables: list[dict[str, Any]] = []

    for i, page in enumerate(reader.pages, start=1):
        raw = page.extract_text() or ""
        page_texts.append(raw)
        if len(raw.strip()) < 40:
            ocr_text, vis, used, n = _ocr_or_vision_pdf_page(page, i)
            ocr_used = ocr_used or used
            notes.extend(n)
            if ocr_text:
                page_texts[-1] = (raw + "\n" + ocr_text).strip()
            visuals.extend(vis)

    # Lightweight markdown tables from pipe/tab-like lines
    combined = "\n\n".join(f"--- Page {i+1} ---\n{t}" for i, t in enumerate(page_texts) if t)
    tables.extend(_infer_text_tables(combined))
    return combined, tables, visuals, ocr_used, notes


def _ocr_or_vision_pdf_page(page, page_no: int) -> tuple[str, list[dict[str, Any]], bool, list[str]]:
    notes = []
    try:
        from pypdf.generic import NameObject
    except Exception:
        pass
    image_bytes = _first_page_image_bytes(page)
    if not image_bytes:
        notes.append(f"Page {page_no} had little selectable text and no embedded image for OCR.")
        return "", [], False, notes
    text, vis, used, extra = _extract_image(image_bytes, f"pdf-page-{page_no}.png")
    for v in vis:
        v["location"] = f"PDF page {page_no}"
    return text, vis, used, notes + extra


def _first_page_image_bytes(page) -> bytes | None:
    try:
        if getattr(page, "images", None):
            img = page.images[0]
            data = getattr(img, "data", None)
            if data:
                return data
    except Exception:
        return None
    return None


def _extract_image(data: bytes, filename: str) -> tuple[str, list[dict[str, Any]], bool, list[str]]:
    notes = []
    ocr_text = _tesseract_ocr(data)
    vision = _gemini_vision(data, filename)
    used = bool(ocr_text or vision)
    parts = []
    visuals = []
    if ocr_text:
        parts.append(ocr_text)
        visuals.append(
            {
                "location": filename,
                "description": ocr_text[:4000],
                "value_kind": "exact",
            }
        )
    if vision:
        parts.append(vision.get("text", ""))
        visuals.append(
            {
                "location": filename,
                "description": vision.get("text", ""),
                "value_kind": vision.get("value_kind", "interpretation"),
            }
        )
    if not used:
        notes.append(
            f"Image '{filename}' could not be OCR'd. Install Tesseract or set GEMINI_API_KEY for vision interpretation."
        )
    return "\n".join(p for p in parts if p), visuals, used, notes


def _tesseract_ocr(data: bytes) -> str:
    try:
        from PIL import Image
        import pytesseract

        img = Image.open(io.BytesIO(data))
        return pytesseract.image_to_string(img) or ""
    except Exception:
        return ""


def _gemini_vision(data: bytes, filename: str) -> dict[str, Any] | None:
    try:
        from rag_chatbot import config
        if not config.GEMINI_API_KEY:
            return None
        from google import genai
        from google.genai import types

        client = genai.Client(api_key=config.GEMINI_API_KEY)
        prompt = (
            "This is a project document image (chart, diagram, screenshot, or scanned page). "
            "Extract all readable text. If it is a graph or chart, describe trend direction, "
            "axis labels, categories, and only approximate values that are actually readable. "
            "Label estimated numbers as ESTIMATED. Do not invent values. "
            "If nothing project-related is readable, say so."
        )
        ext = os.path.splitext(filename)[1].lower()
        mime = MIME_BY_EXT.get(ext, "image/png")
        response = client.models.generate_content(
            model=config.GEMINI_LLM_MODEL,
            contents=[
                types.Part.from_bytes(data=data, mime_type=mime),
                prompt,
            ],
        )
        text = (response.text or "").strip()
        kind = "estimated" if "ESTIMATED" in text.upper() else "interpretation"
        if text:
            return {"text": text, "value_kind": kind}
    except Exception:
        return None
    return None


def _infer_text_tables(text: str) -> list[dict[str, Any]]:
    tables = []
    lines = [ln for ln in text.splitlines() if ln.strip()]
    block: list[str] = []

    def flush():
        if len(block) >= 2:
            rows = [re.split(r"\s{2,}|\t|\|", ln.strip(" |")) for ln in block]
            rows = [[c.strip() for c in r if c.strip()] for r in rows]
            if rows and max(len(r) for r in rows) >= 2:
                headers = rows[0]
                body = rows[1:]
                md = _table_to_markdown(headers, body)
                tables.append({"headers": headers, "rows": body, "markdown": md, "source": "inferred-text-table"})
        block.clear()

    for ln in lines:
        if re.search(r"\|", ln) or re.search(r"\t", ln) or re.search(r"\s{3,}", ln):
            block.append(ln)
        else:
            flush()
    flush()
    return tables[:20]
