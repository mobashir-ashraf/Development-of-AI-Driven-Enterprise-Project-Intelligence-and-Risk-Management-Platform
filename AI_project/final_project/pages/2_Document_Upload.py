import streamlit as st
import pandas as pd
from utils.ui import inject_css, page_header
from utils.predictor import predict_it_risk
from utils.llm_parser import parse_document_with_gemini, parse_batch_with_gemini
from rag_chatbot.session_store import clear_index
from utils.ui import render_model_quality, render_risk_management_processes
from utils.app_store import save_document
from utils.vision_parser import process_document_images

# ============================================================
# PAGE SETUP
# ============================================================

inject_css()

page_header(
    "AI Project Document Upload",
    "Upload your IT project document. The system will autonomously extract scope, tasks, and parameters, while evaluating the project risk."
)

if st.session_state.get("user_type") != "IT":
    st.error("This page is available only for IT users.")
    st.stop()

# ============================================================
# SESSION STATE
# ============================================================

if "documents" not in st.session_state:
    st.session_state.documents = {}
if "project_analyzed" not in st.session_state:
    st.session_state.project_analyzed = False
if "prediction" not in st.session_state:
    st.session_state.prediction = None
if "selected_project_id" not in st.session_state:
    st.session_state.selected_project_id = None
if "selected_project" not in st.session_state:
    st.session_state.selected_project = None
if "it_project_uploaded" not in st.session_state:
    st.session_state.it_project_uploaded = False

# ============================================================
# UTILS
# ============================================================

def extract_text(file):
    filename = file.name.lower()
    if filename.endswith(".txt") or filename.endswith(".csv"):
        try:
            return file.getvalue().decode("utf-8", errors="ignore")
        except: return ""
    if filename.endswith(".docx"):
        try:
            import docx2txt
            text = docx2txt.process(file)
            return text
        except Exception as e:
            st.error(f"Failed to read DOCX: {e}")
            return ""
    if filename.endswith(".pdf"):
        try:
            from pypdf import PdfReader
            reader = PdfReader(file)
            content = [page.extract_text() or "" for page in reader.pages]
            return "\n".join(content)
        except Exception as e:
            st.error(f"Failed to read PDF: {e}")
            return ""
    if filename.endswith(".pptx"):
        try:
            import io
            from pptx import Presentation
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from PIL import Image
            from utils.vision_parser import llm_read_image
            prs = Presentation(io.BytesIO(file.getvalue()))
            text_blocks = []
            for i, slide in enumerate(prs.slides):
                text_blocks.append(f"## Slide {i+1}")
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        text_blocks.append(shape.text)
                    elif shape.has_chart or shape.shape_type == MSO_SHAPE_TYPE.CHART:
                        try:
                            chart_data = []
                            for series in shape.chart.series:
                                series_name = series.name
                                values = series.values
                                categories = [c.label for c in shape.chart.plots[0].categories] if shape.chart.plots else []
                                chart_data.append(f"Series: {series_name}, Categories: {categories}, Values: {values}")
                            text_blocks.append(f"[Chart Data: {chart_data}]")
                        except Exception: pass
                    elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            image_bytes = shape.image.blob
                            img = Image.open(io.BytesIO(image_bytes))
                            if img.mode != "RGB": img = img.convert("RGB")
                            vision_text = llm_read_image(img)
                            if vision_text: text_blocks.append(f"[Image Data: {vision_text}]")
                        except Exception: pass
            return "\n".join(text_blocks)
        except Exception as e:
            st.error(f"Failed to read PPTX: {e}")
            return ""

def save_prediction(insights_dict, prediction):
    project_id = 1000 + len(st.session_state.documents)
    features = insights_dict.get("features", {})
    
    project_data = {
        "id": project_id,
        "project_id": str(project_id),
        "name": insights_dict.get("project_name", "AI Parsed IT Project"),
        "status": "IN_PROGRESS",
        "budget": features.get("budget_usd", 0),
        "deadline": "TBD",
        "progress": 0,
        "health_score": 100 - prediction["risk_score"],
        "risk_level": prediction["risk_level"],
        "risk_score": prediction["risk_score"],
        "features": features,
        
        # New Gemini Insights for UI
        "project_scope": insights_dict.get("project_scope", ""),
        "deliverables": insights_dict.get("deliverables", []),
        "action_items": insights_dict.get("action_items", []),
        "milestones": insights_dict.get("milestones", []),
        "dependencies": insights_dict.get("dependencies", []),
        "missing_info": insights_dict.get("missing_info", []),
        "potential_risks": insights_dict.get("potential_risks", []),
    }
    st.session_state.selected_project_id = project_id
    st.session_state.selected_project = project_data
    st.session_state.project_analyzed = True
    st.session_state.it_project_uploaded = True
    st.session_state.prediction = prediction["risk_level"]

# ============================================================
# UPLOAD
# ============================================================

st.write("### Autonomous Document Processing")
st.write("Upload a PDF, DOCX, CSV, or TXT. Analysis runs locally first, so it remains available without an internet connection.")
render_risk_management_processes()

uploaded_files = st.file_uploader(
    "Choose your IT project document(s)",
    type=["pdf", "docx", "txt", "csv", "png", "jpg", "jpeg", "pptx"],
    accept_multiple_files=True,
    help="Supported formats: PDF, DOCX, TXT, CSV, PPTX, and Images (PNG, JPG, JPEG)"
)

if uploaded_files:
    if st.button("Process Document(s)", type="primary", use_container_width=True):
        st.success(f"Uploaded {len(uploaded_files)} file(s).")
        st.divider()

        is_csv_batch = any(f.name.lower().endswith(".csv") for f in uploaded_files)
        combined_text = ""

        with st.spinner("Reading documents, extracting text and images..."):
            for uploaded_file in uploaded_files:
                file_bytes = uploaded_file.getvalue()
                
                if uploaded_file.name.lower().endswith((".png", ".jpg", ".jpeg")):
                    import io
                    from PIL import Image
                    from utils.vision_parser import llm_read_image
                    img = Image.open(io.BytesIO(file_bytes))
                    if img.mode != "RGB": img = img.convert("RGB")
                    vision_text = llm_read_image(img)
                    extracted_text = f"--- Image Data: {uploaded_file.name} ---\n{vision_text}" if vision_text else ""
                else:
                    extracted_text = extract_text(uploaded_file)
                    
                    # Extract image data using Vision LLM for PDFs
                    if uploaded_file.name.lower().endswith(".pdf"):
                        vision_text = process_document_images(uploaded_file.name, file_bytes)
                        if vision_text:
                            extracted_text += "\n" + vision_text

                if not extracted_text.strip():
                    continue

                st.session_state.documents[uploaded_file.name] = extracted_text
                combined_text += f"\n\n--- Document: {uploaded_file.name} ---\n{extracted_text}"

                # --- Compute per-file-type stats for Project Analysis ---
                fname_lower = uploaded_file.name.lower()
                doc_meta: dict = {"size_bytes": len(file_bytes)}
                if fname_lower.endswith(".pdf"):
                    doc_meta["doc_type"] = "PDF"
                    try:
                        from pypdf import PdfReader
                        import io as _io
                        _rdr = PdfReader(_io.BytesIO(file_bytes))
                        doc_meta["page_count"] = len(_rdr.pages)
                    except Exception:
                        doc_meta["page_count"] = None
                elif fname_lower.endswith(".pptx"):
                    doc_meta["doc_type"] = "PPTX"
                    try:
                        import io as _io
                        from pptx import Presentation as _Prs
                        _prs = _Prs(_io.BytesIO(file_bytes))
                        doc_meta["slide_count"] = len(_prs.slides)
                    except Exception:
                        doc_meta["slide_count"] = None
                elif fname_lower.endswith(".docx"):
                    doc_meta["doc_type"] = "DOCX"
                    words = len(extracted_text.split())
                    doc_meta["page_count"] = max(1, round(words / 250))
                elif fname_lower.endswith(".csv"):
                    doc_meta["doc_type"] = "CSV"
                    doc_meta["row_count"] = max(0, extracted_text.count("\n") - 1)
                elif fname_lower.endswith(".txt"):
                    doc_meta["doc_type"] = "TXT"
                    doc_meta["char_count"] = len(extracted_text)
                elif fname_lower.endswith((".png", ".jpg", ".jpeg")):
                    doc_meta["doc_type"] = "Image"
                    doc_meta["image_count"] = 1
                else:
                    doc_meta["doc_type"] = "Document"

                # Save to persistent database
                save_document(
                    user_id=st.session_state.user_id,
                    filename=uploaded_file.name,
                    content=file_bytes,
                    extracted_text=extracted_text,
                    metadata=doc_meta,
                )

            if not combined_text.strip():
                st.error("No readable text was found. Upload text-based or valid PDF/DOCX files.")
                st.stop()
            
            try:
                if is_csv_batch:
                    insights_dict = parse_batch_with_gemini(combined_text)
                else:
                    insights_dict = parse_document_with_gemini(combined_text)
            except Exception as e:
                st.error(f"Failed to process and analyze document: {e}")
                st.stop()
                
        with st.spinner("Running XGBoost Risk Model..."):
            if is_csv_batch:
                batch_projects = []
                for i, proj in enumerate(insights_dict.get("projects", [])):
                    features = proj.get("features", {})
                    prediction = predict_it_risk(features)
                    
                    project_id = 2000 + i
                    project_data = {
                        "id": project_id,
                        "project_id": str(project_id),
                        "name": proj.get("project_name", f"Batch Project {i+1}"),
                        "status": "IN_PROGRESS",
                        "budget": features.get("budget_usd", 0),
                        "deadline": "TBD",
                        "progress": 0,
                        "health_score": 100 - prediction["risk_score"],
                        "risk_level": prediction["risk_level"],
                        "risk_score": prediction["risk_score"],
                        "features": features,
                        "project_scope": proj.get("project_scope", ""),
                        "deliverables": proj.get("deliverables", []),
                        "action_items": proj.get("action_items", []),
                        "milestones": proj.get("milestones", []),
                        "dependencies": proj.get("dependencies", []),
                        "missing_info": proj.get("missing_info", []),
                        "potential_risks": proj.get("potential_risks", []),
                    }
                    batch_projects.append(project_data)
                    
                st.session_state.batch_projects = batch_projects
                
                # Default to the first project in the batch for deep dives
                if batch_projects:
                    st.session_state.selected_project = batch_projects[0]
                    st.session_state.selected_project_id = batch_projects[0]["id"]
                    
                st.session_state.project_analyzed = False
                st.session_state.is_batch = True
            else:
                features = insights_dict.get("features", {})
                prediction = predict_it_risk(features)
                save_prediction(insights_dict, prediction)
                st.session_state.is_batch = False

        # Indexing is intentionally deferred until the user opens RAGBot.  It
        # can be expensive for large files and must not delay core analysis.
        clear_index()
        st.session_state["rag_ready"] = False
        st.session_state["rag_chunk_count"] = 0

        st.success("Project analyzed autonomously! Check the Dashboards for details.")
        st.info("RAGBot indexing is available on demand from the AI Assistant page.")

# ============================================================
# DISPLAY ANALYSIS
# ============================================================

if st.session_state.get("is_batch", False) and "batch_projects" in st.session_state:
    st.divider()
    st.success("Batch processing complete!")
    st.write(f"Successfully analyzed {len(st.session_state.batch_projects)} projects.")
    st.page_link("pages/1_Dashboard.py", label="View Batch Risk Distribution Dashboard")

elif st.session_state.get("project_analyzed", False):
    project = st.session_state.get("selected_project")
    if not project: st.stop()

    st.divider()
    st.subheader("Project Analysis Snapshot")

    col1, col2, col3, col4 = st.columns(4)
    with col1: st.metric("Project Name", project.get("name"))
    with col2: st.metric("Overall Health", f"{project.get('health_score')}/100")
    with col3: st.metric("XGBoost Risk Score", f"{project.get('risk_score')}/100")
    with col4: st.metric("AI Risk Level", project.get("risk_level"))

    st.markdown(f"""
    <div class="glass-card" style="margin-top: 1rem; margin-bottom: 1.5rem; border-color: rgba(56, 189, 248, 0.35);">
        <div style="display: flex; align-items: center; gap: 0.6rem; margin-bottom: 0.6rem;">
            <span style="font-size: 1.3rem;"></span>            <span style="color: #38bdf8; font-weight: 800; font-size: 1.05rem; text-transform: uppercase; letter-spacing: 0.5px;">Executive Document Scope & Highlights</span>
        </div>
        <p style="color: #e2e8f0; font-size: 0.96rem; line-height: 1.6; margin: 0;">{project.get('project_scope', 'Extracted scope details from uploaded project documentation.')}</p>
    </div>
    """, unsafe_allow_html=True)

    details_col, risk_col = st.columns(2)
    with details_col:
        with st.expander("Strategic Deliverables & Milestone Telemetry", expanded=True):
            st.markdown("<h4 style='color: #ffffff; font-size: 0.95rem; margin-bottom: 0.6rem;'>Key Engineering Workstreams</h4>", unsafe_allow_html=True)
            for item in project.get("deliverables", [])[:6]:
                st.markdown(f"<div style='color: #cbd5e1; font-size: 0.88rem; margin-bottom: 0.35rem;'><strong>{item}</strong></div>", unsafe_allow_html=True)
            
            st.divider()
            st.markdown("<h4 style='color: #ffffff; font-size: 0.95rem; margin-bottom: 0.6rem;'>Milestone Completion Telemetry</h4>", unsafe_allow_html=True)
            for milestone in project.get("milestones", [])[:4]:
                name = milestone.get('name', 'Unnamed Phase')
                pct = int(milestone.get('progress_pct', 0))
                st.caption(f"{name} — **{pct}% Complete**")
                st.progress(pct)

    with risk_col:
        with st.expander("Predictive Risk Telemetry & Evidence", expanded=True):
            features = project.get("features", {})
            evidence = {
                "Capital Budget Baseline": f"${features.get('budget_usd', 0):,.0f}",
                "Schedule Variance Exposure": f"{features.get('schedule_overrun_pct', 0):.1f}%",
                "Resource Availability Index": f"{features.get('resource_availability_pct', 0):.0f}%",
                "Technical Complexity Rating": f"{features.get('tech_complexity_score', 0):.0f}/100",
                "External Dependency Index": f"{features.get('external_dependency_score', 0):.0f}/100",
            }
            st.dataframe(pd.DataFrame(evidence.items(), columns=["Telemetry Signal", "Extracted Value"]), hide_index=True, use_container_width=True)
            st.markdown("<h4 style='color: #ffffff; font-size: 1.05rem; margin-top: 1rem; margin-bottom: 0.6rem; font-weight: 800;'>Key Risk Triggers</h4>", unsafe_allow_html=True)
            for risk in project.get("potential_risks", [])[:5]:
                r_clean = str(risk).strip(" -•\t")
                if r_clean:
                    r_cap = r_clean[0].upper() + r_clean[1:]
                    st.markdown(f"<div style='color: #f8fafc; font-size: 0.98rem; font-weight: 500; margin-bottom: 0.45rem; line-height: 1.5;'>{r_cap}</div>", unsafe_allow_html=True)
    
    st.write("Use the sidebar navigation to see deep dives into Action Items, Dashboards, and generated documentation.")
    render_model_quality("IT")
